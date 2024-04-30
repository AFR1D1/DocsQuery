import sys
import subprocess
from google.colab import files
import tempfile
import shutil
import os
import time



library_names = ['spacy', 'spacy[transformers]','pytesseract', 'sentence-transformers', 'langchain', 'langchain-openai', 'faiss-cpu', 'PyPDF2','python-docx', 'openai', 'tiktoken', 'python-pptx', 'textwrap', ]

# Dynamically importing libraries
for name in library_names:
    try:
        __import__(name)
    except ImportError:
        print(f"{name} not found. Installing {name}...")
        subprocess.check_call([sys.executable, '-m', 'pip', 'install', name])

# Install spacy[transformers] and download en_core_web_trf if spacy is successfully imported
try:
    __import__('spacy')
    subprocess.check_call([sys.executable, '-m', 'spacy', 'download', 'en_core_web_trf'])
except Exception as e:
    print(f"An error occurred: {e}")

from pytesseract import image_to_string
from PIL import Image
from PyPDF2 import PdfReader 
import textwrap
import docx
import pptx
from langchain_openai import OpenAIEmbeddings, OpenAI
from langchain.text_splitter import CharacterTextSplitter
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from getpass import getpass
import io

# Import spaCy after installation
import spacy
from spacy.matcher import Matcher


import torch
from sentence_transformers import SentenceTransformer, util

# Downloading the English language model for spaCy
#try:
nlp = spacy.load("en_core_web_trf")
#except IOError:
#    print("en_core_web_trf not found. Attempting to download...")
#    subprocess.check_call([sys.executable, '-m', 'spacy', 'download', 'en_core_web_trf'])
#    nlp = spacy.load("en_core_web_trf")

matcher = Matcher(nlp.vocab)

# Add match pattern for URLs
pattern = [{"LIKE_URL": True}]
matcher.add("URL_PATTERN", [pattern])

model = SentenceTransformer('all-MiniLM-L6-v2')

def extract_keywords(text):
    doc = nlp(text)
    keywords = set()
    
    # Extract entities and nouns as keywords
    for ent in doc.ents:
        keywords.add(ent.lemma_.lower())
    for token in doc:
        if token.pos_ in ['NOUN', 'PROPN'] and not token.is_stop:
            keywords.add(token.lemma_.lower())
    
    return list(keywords)




#token adding
if "OPENAI_API_KEY" in os.environ:
    print("Token already set.")
else:
    token = getpass("Enter your OpenAI token: ")
    os.environ["OPENAI_API_KEY"] = str(token)





# Downloading embeddings from OpenAI
embeddings = OpenAIEmbeddings()
chain = load_qa_chain(OpenAI(), chain_type="stuff")


def extract_texts(root_files):
    """
    Text extractes from file and puts it in a list.
    Supported file formats include: .pdf, .docx, .pptx
    If multiple files are uploaded, contents will be merged together
    Parameters:
    - root_files: A list containing the paths of the files to be processed.
    Returns:
    - A FAISS index object that includes the embeddings of the extracted text segments.
    """
    raw_text = ''

    for root_file in root_files:
        _, ext = os.path.splitext(root_file)
        if ext == '.pdf':
            # First try to extract text normally
            with open(root_file, 'rb') as f:
                reader = PdfReader(f)
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        raw_text += text + '\n'
                    else:
                        # If normal text extraction doesn't work, use OCR
                        images = convert_from_path(root_file)
                        for image in images:
                            raw_text += image_to_string(image) + '\n'
        elif ext == '.docx':
            doc = docx.Document(root_file)
            for paragraph in doc.paragraphs:
                raw_text += paragraph.text + '\n'
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    image_stream = io.BytesIO(rel.target_part.blob)
                    image = Image.open(image_stream)
                    raw_text += image_to_string(image) + '\n'
        elif ext == '.pptx':
            ppt = pptx.Presentation(root_file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        raw_text += shape.text + '\n'
                    elif shape.shape_type == 13: # ShapeType 13 corresponds to a picture
                        image_stream = io.BytesIO(shape.image.blob)
                        image = Image.open(image_stream)
                        raw_text += image_to_string(image) + '\n'

    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
    )

    texts = text_splitter.split_text(raw_text)
    docsearch = FAISS.from_texts(texts, embeddings)

    # Return both the FAISS index and the texts
    return docsearch, texts


def run_query(query, docsearch):
    """
    Executes a search query on a PDF file utilizing the docsearch and chain libraries.
    Parameters:
    query: A string that specifies the query to be executed.
    file: A PDFReader object that holds the PDF file to be queried.
    Returns:
    A string that includes the results from applying the chain library to the documents retrieved by the docsearch similarity search.
    """
    
    docs = docsearch.similarity_search(query)
    return chain.run(input_documents=docs, question=query)


def upload_file(folder_path):
    """
    Uploads a file from the local file system and stores it in a specified directory.
    Parameters:
    folder_path: A string that indicates the directory where the file should be saved.
    Returns:
    A string that denotes the path of the file that has been uploaded.
    """
    
    uploaded = files.upload()
    root_file = []

    for filename, data in uploaded.items():
        with open(filename, 'wb') as f:
            f.write(data)
        shutil.copy(filename, folder_path + "/")
        root_file.append(folder_path + "/" + filename)
        os.remove(filename)


    return root_file
################











##############


def find_similar_questions(query, questions, top_k=5):
    query_embedding = model.encode(query, convert_to_tensor=True)
    question_embeddings = model.encode(questions, convert_to_tensor=True)
    cos_scores = util.pytorch_cos_sim(query_embedding, question_embeddings)[0]
    top_results = torch.topk(cos_scores, k=top_k)
    return [questions[index] for index in top_results.indices]

def load_all_questions(all_texts):
    """
    Dynamically generates questions based on the content of uploaded files.
    """
    unique_keywords = set()
    for text in all_texts:
        keywords = extract_keywords(text)
        unique_keywords.update(keywords)
    
    # Generate questions from keywords
    questions = []
    for keyword in unique_keywords:
        if len(keyword) > 2:  # Filter out very short keywords that might be less relevant
            questions.extend([
                f"What is {keyword}?",
                f"How does {keyword} work?",
                f"What are the applications of {keyword}?",
                f"Explain the concept of {keyword}",
                f"Advantages and disadvantages of {keyword}?",
                f"What theoretical frameworks underpin {keyword}?",
                f"How has {keyword} evolved over time?",
                f"Compare {keyword} with another concept in its category.",
                f"What methodologies are most effective for studying {keyword}?",
                f"What ethical issues surround {keyword}?",
                f"What are the future research directions for {keyword}?"
            ])
    return questions


#################

def run_conversation(folder_path):
    root_files = upload_file(folder_path)
    docsearch, all_texts = extract_texts(root_files)  # Receive texts as well
    all_questions = load_all_questions(all_texts)
    count = 0
    while True:
        print(f"Question {count + 1}")
        user_input = input("Ask questions, type suggest, or type stop:\n")
        if user_input.lower() == "stop":
            print("Thanks for using the service.")
            break
        elif user_input.lower() == "suggest":
            # Generate and print related questions suggestions
            if count == 0:
                print("Please ask a question first before suggesting related topics.")
            else:
                suggestions = find_similar_questions(last_query, all_questions)
                print("Related questions:")
                for question in suggestions:
                    print(question)
        elif user_input.strip() == "":
            print("Input is empty, please enter a valid command.")
        else:
            # Extract keywords from the user's query
            keywords = extract_keywords(user_input)
            
            # Run the query against the documents
            response = run_query(user_input, docsearch)
            
            # Print the response
            wrapped_text = textwrap.wrap(response, width=100)
            print("Answer:")
            for line in wrapped_text:
                print(line)
            
            # Save the last valid question for the suggest feature
            last_query = user_input
            
            count += 1
